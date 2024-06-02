from VESTER_AI import app
from flask import request, jsonify
import os
from PyPDF2 import PdfReader, DocumentInformation
import pprint
from pptx import Presentation


# Set the directory where uploaded files will be saved
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Ensure the upload folder exists
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)


# Allowed extensions
ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx'}
ALLOWED_MIME_TYPES = {'application/pdf', 'application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', "application/octet-stream"}

def allowed_file(filename, file):
    # Check file extension
    extension = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    if extension not in ALLOWED_EXTENSIONS:
        return False

    # Check MIME type
    if file.mimetype not in ALLOWED_MIME_TYPES:
        return False

    return True


@app.route('/')
def index():
    return {
        "statusCode": 200,
        "message": "nothing to see here!"
    }


@app.route('/upload', methods=['POST'])
def upload_file():
    #  [ VALIDATE REQUEST - Check if the post request has the file part ]
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
    
    file = request.files['file']

    """
    If user does not select file, browser also
    submits an empty part without filename
    """
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if file and allowed_file(file.filename, file):
        # Save the file to the specified folder
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)

        if file.mimetype == "application/pdf":
            # to handle for PDFs only
            try:
                # for pdf
                reader = PdfReader(file.filename)
                meta = reader.metadata

                number_of_pages = len(reader.pages)
                texts = {}

                for i in range(1, number_of_pages):
                    page = reader.pages[i]
                    text = page.extract_text()
                    texts[i] = text
            except Exception as e:
                return jsonify({
                        "statusCode": 400,
                        "err": str(e)
                })

            pprint.pprint(texts) # Will save to db
            return jsonify({
                "statusCode": 200,
                "message": "File successfully uploaded",
                "working_dir": os.getcwd(),
                "file_path": file_path,
                "file_name": file.filename,
                "file_type": file.mimetype,
                "file"
                "reader": {
                    "title": meta.title,
                    "author": meta.author,
                    "number_of_pages": number_of_pages
                }
            })
        elif file.mimetype in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/octet-stream"
        ]:
            # Too handle for powerpoint files.
            _path = os.getcwd() + "/" + file_path
            # return _path
            prs = Presentation(_path)

            for slide_number, slide in enumerate(prs.slides):
                print(f"Slide {slide_number + 1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        print(shape.text)

            return {
                "statusCode": 200,
                "file_name": file.filename,
                "file_path": _path,
                "file_type": file.mimetype
            }
    else:
        return jsonify({"error": "Invalid file type. Only PDF and PowerPoint files are allowed."}), 400

